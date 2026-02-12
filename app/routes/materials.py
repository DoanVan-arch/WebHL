from fastapi import APIRouter, Depends, HTTPException, status, UploadFile, File, Form
from sqlalchemy.orm import Session
from typing import List, Optional
from app.database.database import get_db
from app.models.models import Material, Department, User, UserRole
from app.dependencies import get_current_user
import os
import shutil
import json
from datetime import datetime
from docx import Document
from pptx import Presentation
from PyPDF2 import PdfReader

router = APIRouter()

UPLOAD_DIR = "app/static/uploads"
os.makedirs(UPLOAD_DIR, exist_ok=True)

def extract_text_from_file(file_path: str) -> str:
    """Extract text content from PDF, DOCX, or PPTX files"""
    try:
        ext = os.path.splitext(file_path)[1].lower()
        
        if ext == '.pdf':
            reader = PdfReader(file_path)
            text = ""
            for page in reader.pages:
                text += page.extract_text() + "\n"
            return text
        
        elif ext == '.docx':
            doc = Document(file_path)
            text = "\n".join([paragraph.text for paragraph in doc.paragraphs])
            return text
        
        elif ext == '.pptx':
            prs = Presentation(file_path)
            text = ""
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
            return text
        
        return ""
    except Exception as e:
        print(f"Error extracting text from {file_path}: {e}")
        return ""

@router.get("/materials")
async def get_materials(
    department_id: Optional[int] = None,
    search: Optional[str] = None,
    uploader: Optional[str] = None,
    search_content: Optional[str] = None,
    db: Session = Depends(get_db),
    current_user: User = Depends(get_current_user)
):
    query = db.query(Material)
    
    if department_id:
        query = query.filter(Material.department_id == department_id)
    
    if uploader:
        # Tìm kiếm theo tên người đăng
        uploader_filter = f"%{uploader}%"
        query = query.join(User).filter(User.full_name.like(uploader_filter))
    
    materials = query.order_by(Material.created_at.desc()).all()
    
    # Tìm kiếm theo nội dung file nếu được bật
    if search and search_content == 'true':
        filtered_materials = []
        search_lower = search.lower()
        
        for m in materials:
            # Kiểm tra metadata trước
            if (search_lower in m.title.lower() or 
                search_lower in m.subject.lower() or 
                (m.topic and search_lower in m.topic.lower())):
                filtered_materials.append(m)
                continue
            
            # Tìm trong nội dung file
            files = json.loads(m.files_json)
            found = False
            for file_info in files:
                file_path = file_info['path'].replace('/static/', 'app/static/')
                if os.path.exists(file_path):
                    content = extract_text_from_file(file_path)
                    if search_lower in content.lower():
                        filtered_materials.append(m)
                        found = True
                        break
            
        materials = filtered_materials
    elif search:
        # Tìm kiếm thông thường theo metadata
        search_filter = f"%{search}%"
        materials = [m for m in materials if (
            search.lower() in m.title.lower() or
            search.lower() in m.subject.lower() or
            (m.topic and search.lower() in m.topic.lower())
        )]
    
    return {
        "materials": [
            {
                "id": m.id,
                "title": m.title,
                "subject": m.subject,
                "topic": m.topic,
                "files": json.loads(m.files_json),
                "department": {
                    "id": m.department.id,
                    "code": m.department.code,
                    "name": m.department.name
                },
                "uploader": {
                    "id": m.uploader.id,
                    "full_name": m.uploader.full_name
                },
                "created_at": m.created_at.isoformat()
            }
            for m in materials
        ]
    }

@router.get("/materials/{material_id}")
async def get_material_detail(
    material_id: int,
    db: Session = Depends(get_db),
    current_user: User = Depends(get_current_user)
):
    material = db.query(Material).filter(Material.id == material_id).first()
    
    if not material:
        raise HTTPException(
            status_code=status.HTTP_404_NOT_FOUND,
            detail="Không tìm thấy học liệu"
        )
    
    return {
        "id": material.id,
        "title": material.title,
        "subject": material.subject,
        "topic": material.topic,
        "files": json.loads(material.files_json),
        "department": {
            "id": material.department.id,
            "code": material.department.code,
            "name": material.department.name
        },
        "uploader": {
            "id": material.uploader.id,
            "full_name": material.uploader.full_name,
            "email": material.uploader.email
        },
        "created_at": material.created_at.isoformat(),
        "updated_at": material.updated_at.isoformat()
    }

@router.post("/materials")
async def create_material(
    title: str = Form(...),
    subject: str = Form(...),
    topic: str = Form(""),
    department_id: int = Form(...),
    tailieu_files: List[UploadFile] = File(default=[]),
    baigiang_files: List[UploadFile] = File(default=[]),
    decuong_files: List[UploadFile] = File(default=[]),
    trinhchieu_files: List[UploadFile] = File(default=[]),
    db: Session = Depends(get_db),
    current_user: User = Depends(get_current_user)
):
    # Check if user is superuser or admin
    if current_user.role not in [UserRole.SUPERUSER, UserRole.ADMIN]:
        raise HTTPException(
            status_code=status.HTTP_403_FORBIDDEN,
            detail="Chỉ superuser và admin mới có quyền đăng học liệu"
        )
    
    # Collect all files
    all_files = []
    file_groups = [
        ("Tài liệu", tailieu_files),
        ("Bài giảng", baigiang_files),
        ("Đề cương", decuong_files),
        ("Trình chiếu", trinhchieu_files)
    ]
    
    # Check if at least one file is uploaded
    total_files = sum(len(files) for _, files in file_groups if files)
    if total_files == 0:
        raise HTTPException(
            status_code=status.HTTP_400_BAD_REQUEST,
            detail="Vui lòng chọn ít nhất một file"
        )
    
    # Save files
    for file_type, files in file_groups:
        if files:
            for file in files:
                if file.filename:
                    timestamp = datetime.now().strftime("%Y%m%d_%H%M%S_%f")
                    file_extension = os.path.splitext(file.filename)[1]
                    unique_filename = f"{timestamp}_{file.filename}"
                    file_path = os.path.join(UPLOAD_DIR, unique_filename)
                    
                    with open(file_path, "wb") as buffer:
                        shutil.copyfileobj(file.file, buffer)
                    
                    all_files.append({
                        "type": file_type,
                        "path": f"/static/uploads/{unique_filename}",
                        "name": file.filename
                    })
    
    # Create material record
    new_material = Material(
        title=title,
        subject=subject,
        topic=topic,
        files_json=json.dumps(all_files, ensure_ascii=False),
        department_id=department_id,
        uploader_id=current_user.id
    )
    
    db.add(new_material)
    db.commit()
    db.refresh(new_material)
    
    return {
        "message": "Đăng học liệu thành công",
        "material_id": new_material.id
    }

@router.delete("/materials/{material_id}")
async def delete_material(
    material_id: int,
    db: Session = Depends(get_db),
    current_user: User = Depends(get_current_user)
):
    material = db.query(Material).filter(Material.id == material_id).first()
    
    if not material:
        raise HTTPException(
            status_code=status.HTTP_404_NOT_FOUND,
            detail="Không tìm thấy học liệu"
        )
    
    # Check permissions
    if current_user.role == UserRole.USER:
        raise HTTPException(
            status_code=status.HTTP_403_FORBIDDEN,
            detail="Bạn không có quyền xóa học liệu"
        )
    
    if current_user.role == UserRole.SUPERUSER and material.uploader_id != current_user.id:
        raise HTTPException(
            status_code=status.HTTP_403_FORBIDDEN,
            detail="Bạn chỉ có thể xóa học liệu của mình"
        )
    
    # Delete files
    files = json.loads(material.files_json)
    for file_info in files:
        file_path = f"app{file_info['path']}"
        if os.path.exists(file_path):
            os.remove(file_path)
    
    db.delete(material)
    db.commit()
    
    return {"message": "Xóa học liệu thành công"}

@router.put("/materials/{material_id}")
async def update_material(
    material_id: int,
    title: str = Form(...),
    subject: str = Form(...),
    topic: str = Form(None),
    department_id: int = Form(...),
    db: Session = Depends(get_db),
    current_user: User = Depends(get_current_user)
):
    """Update material information (metadata only, not files)"""
    material = db.query(Material).filter(Material.id == material_id).first()
    
    if not material:
        raise HTTPException(
            status_code=status.HTTP_404_NOT_FOUND,
            detail="Không tìm thấy học liệu"
        )
    
    # Check permissions
    if current_user.role == UserRole.USER:
        raise HTTPException(
            status_code=status.HTTP_403_FORBIDDEN,
            detail="Bạn không có quyền chỉnh sửa học liệu"
        )
    
    if current_user.role == UserRole.SUPERUSER and material.uploader_id != current_user.id:
        raise HTTPException(
            status_code=status.HTTP_403_FORBIDDEN,
            detail="Bạn chỉ có thể chỉnh sửa học liệu của mình"
        )
    
    # Check if department exists
    department = db.query(Department).filter(Department.id == department_id).first()
    if not department:
        raise HTTPException(
            status_code=status.HTTP_404_NOT_FOUND,
            detail="Không tìm thấy khoa"
        )
    
    # Update material
    material.title = title
    material.subject = subject
    material.topic = topic
    material.department_id = department_id
    material.updated_at = datetime.now()
    
    db.commit()
    db.refresh(material)
    
    return {
        "message": "Cập nhật học liệu thành công",
        "material": {
            "id": material.id,
            "title": material.title,
            "subject": material.subject,
            "topic": material.topic,
            "department_id": material.department_id
        }
    }

@router.get("/preview/{material_id}/{file_index}")
async def preview_file(
    material_id: int,
    file_index: int,
    db: Session = Depends(get_db),
    current_user: User = Depends(get_current_user)
):
    """Preview DOCX or PPTX file content (first 5 pages/slides)"""
    material = db.query(Material).filter(Material.id == material_id).first()
    
    if not material:
        raise HTTPException(
            status_code=status.HTTP_404_NOT_FOUND,
            detail="Không tìm thấy học liệu"
        )
    
    files = json.loads(material.files_json)
    
    if file_index >= len(files):
        raise HTTPException(
            status_code=status.HTTP_404_NOT_FOUND,
            detail="Không tìm thấy file"
        )
    
    file_info = files[file_index]
    file_path = file_info['path'].replace('/static/', 'app/static/')
    
    if not os.path.exists(file_path):
        raise HTTPException(
            status_code=status.HTTP_404_NOT_FOUND,
            detail="File không tồn tại"
        )
    
    ext = os.path.splitext(file_path)[1].lower()
    
    try:
        if ext == '.docx':
            print(f"[PREVIEW] Reading DOCX file: {file_path}")
            doc = Document(file_path)
            paragraphs = []
            page_count = 0
            max_pages = 5
            
            # Estimate pages (rough estimate: ~40 paragraphs per page)
            for i, para in enumerate(doc.paragraphs):
                if page_count >= max_pages:
                    break
                paragraphs.append({
                    'text': para.text,
                    'style': para.style.name if para.style else 'Normal'
                })
                if i % 40 == 0 and i > 0:
                    page_count += 1
            
            print(f"[PREVIEW] Successfully read DOCX: {len(paragraphs)} paragraphs")
            return {
                'type': 'docx',
                'total_paragraphs': len(doc.paragraphs),
                'preview_paragraphs': len(paragraphs),
                'paragraphs': paragraphs,
                'estimated_pages': len(doc.paragraphs) // 40 + 1
            }
        
        elif ext == '.pptx':
            print(f"[PREVIEW] Reading PPTX file: {file_path}")
            prs = Presentation(file_path)
            slides_data = []
            
            for i, slide in enumerate(prs.slides):
                if i >= 5:  # Only first 5 slides
                    break
                
                slide_content = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_content.append({
                            'text': shape.text,
                            'type': str(shape.shape_type) if hasattr(shape, 'shape_type') else 'UNKNOWN'
                        })
                
                slides_data.append({
                    'slide_number': i + 1,
                    'content': slide_content
                })
            
            print(f"[PREVIEW] Successfully read PPTX: {len(slides_data)} slides")
            return {
                'type': 'pptx',
                'total_slides': len(prs.slides),
                'preview_slides': len(slides_data),
                'slides': slides_data
            }
        
        else:
            raise HTTPException(
                status_code=status.HTTP_400_BAD_REQUEST,
                detail="Chỉ hỗ trợ preview file DOCX và PPTX"
            )
    
    except HTTPException:
        raise
    except Exception as e:
        print(f"[PREVIEW ERROR] Failed to read file: {file_path}")
        print(f"[PREVIEW ERROR] Error: {str(e)}")
        import traceback
        traceback.print_exc()
        raise HTTPException(
            status_code=status.HTTP_500_INTERNAL_SERVER_ERROR,
            detail=f"Lỗi khi đọc file: {str(e)}"
        )

@router.get("/departments")
async def get_departments(
    db: Session = Depends(get_db),
    current_user: User = Depends(get_current_user)
):
    departments = db.query(Department).all()
    return {
        "departments": [
            {
                "id": d.id,
                "code": d.code,
                "name": d.name,
                "description": d.description
            }
            for d in departments
        ]
    }
